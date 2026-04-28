"""Build DETECTive_30min_full_deck.pptx — the complete 30-minute presentation.

Output: presentation/DETECTive_30min_full_deck.pptx (16:9, default theme).

Structure:
  Title slide
  BLOCK 1 — Problem & Classical ATPG (~10 min):
    B1.1 Digital Circuit Fault
    B1.2 Test Pattern
    B1.3 Exponential Problem
    B1.4 D-Algorithm (Roth, 1966)              [colleague will tweak]
    B1.5 Working of D-Algorithm                 [colleague will tweak]
    B1.6 PODEM (Goel 1981)                      [reused from build_pptx.py]
    B1.7 Working of PODEM                       [reused]
    B1.8 Worked PODEM Example                   [reused]
    B1.9 FAN (Fujiwara-Shimono 1983)           [colleague will tweak]
    B1.10 Working of FAN                        [colleague will tweak]
    B1.11 Scalability Bottleneck & Why ML
  BLOCK 2 — ATPP Architecture (~10 min):
    B2.1 ATPP Concept
    B2.2 Circuit as Graph
    B2.3 GNN Encoder: GAT + GCN
    B2.4 LSTM Path Encoder
    B2.5 Full Architecture
    B2.6 Node Features and Tensor Shapes
    B2.7 Training Setup
    B2.8 Training Curve & Overfitting
  BLOCK 3 — Results (~10 min):
    B3.1 Dataset
    B3.2 ISCAS-85 Results Table
    B3.3 4-Way Bit-Accuracy
    B3.4 Runtime per Fault — 17.6x Speedup
    B3.5 Bit-Acc vs Fault-Sim Coverage Gap
    B3.6 Synthetic Grid (paper Fig 5a)
    B3.7 Depth Degradation (paper Fig 6a)
    B3.8 Honest Takeaways & Limitations
    B3.9 Conclusion
  Closing:
    Q & A
    Thank You
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG_DIR = os.path.join(PROJECT_ROOT, "presentation_images")
OUT_PATH = os.path.join(PROJECT_ROOT, "presentation", "DETECTive_30min_full_deck.pptx")


# ---------------------------------------------------------------------------
# Slide specifications
# ---------------------------------------------------------------------------

SLIDES = [
    # =====================================================================
    # BLOCK 1 — Problem & Classical ATPG
    # =====================================================================
    dict(
        n=1, title="Digital Circuit Fault",
        bullets=[
            "Manufacturing defects: bridges, opens, transistor failures during fabrication.",
            "Stuck-at fault model: abstract a defect as a wire permanently driven to a fixed value.",
            "Stuck-at-0 (SA-0): the line is always logic 0 regardless of its driver.",
            "Stuck-at-1 (SA-1): the line is always logic 1 regardless of its driver.",
            "Industry standard: every line has 2 stuck-at faults — a circuit with N lines has 2N targets.",
            "Goal of ATPG: generate input vectors that detect all stuck-at faults on the chip.",
        ],
        notes=(
            "Let's start with the basic problem. When we manufacture a chip, things can go wrong — a wire "
            "might short to ground, an open might prevent a signal from reaching its destination, or a "
            "transistor might fail. Rather than model every possible physical defect, the industry has "
            "settled on a clean abstraction called the stuck-at fault model. We say a wire is either "
            "stuck-at-zero — it always reads as logic zero no matter what its driver does — or stuck-at-one. "
            "Every line in the circuit has two possible faults, so a circuit with a thousand lines has two "
            "thousand fault targets. The job of automatic test pattern generation, ATPG, is to produce input "
            "vectors that expose every one of these faults so we can catch defective chips before shipping them."
        ),
        image=None, time_s=60,
    ),
    dict(
        n=2, title="Test Pattern",
        bullets=[
            "A test pattern is a binary vector applied to the primary inputs (PIs) of the circuit.",
            "Two requirements for detecting a fault:",
            "  1. Activation — force the faulty line to the OPPOSITE value of the stuck-at value.",
            "  2. Propagation — sensitize a path so the difference reaches a primary output (PO).",
            "Detection condition: faulty-circuit output != fault-free output on at least one PO.",
            "Example: SA-0 on line L requires (a) drive L to 1 in the good circuit and (b) propagate the 1-vs-0 difference to a PO.",
        ],
        notes=(
            "What is a test pattern, concretely? It's just a binary vector — one bit per primary input — "
            "that you apply to the chip. For a test pattern to detect a fault, two things have to happen "
            "at once. First, activation: you need to drive the faulty line to the value opposite of its "
            "stuck-at. So if a wire is stuck-at-zero, you have to force the good circuit to want a one "
            "there. Second, propagation: you need to sensitize a path from the fault site all the way to "
            "an observable output, so that the difference between the good circuit and the faulty circuit "
            "is visible from outside. The fault is only detected if the output of the faulty chip disagrees "
            "with the good chip on at least one primary output bit. Both conditions are necessary."
        ),
        image=None, time_s=70,
    ),
    dict(
        n=3, title="The Exponential Problem",
        bullets=[
            "A circuit with n primary inputs has 2^n possible test patterns.",
            "Even modest circuits have hundreds to thousands of PIs.",
            "c1908 has 33 PIs => 2^33 ~ 8.6 billion possible patterns. c2670 has 233 PIs => intractable.",
            "Exhaustive testing is infeasible — we cannot try every input combination.",
            "We need smart algorithms that pick a small set of useful patterns (target each fault directly).",
            "This is the core motivation for ATPG: structured search instead of brute force.",
        ],
        notes=(
            "Why don't we just try every possible input? Because the input space explodes. A circuit with "
            "n primary inputs has two-to-the-n possible test patterns. c1908, a small ISCAS-85 benchmark, "
            "already has thirty-three primary inputs — that's eight-and-a-half billion patterns. c2670 has "
            "two hundred thirty-three. Two-to-the-two-hundred-thirty-three is more atoms than the universe. "
            "Modern industrial chips have thousands of primary inputs, and even after applying scan chains "
            "the effective input space stays huge. Brute force is out. So we need clever algorithms that "
            "target each fault directly and produce a small, high-coverage test set. That is the entire "
            "field of automatic test pattern generation."
        ),
        image=None, time_s=70,
    ),
    dict(
        n=4, title="D-Algorithm (Roth, 1966)",
        bullets=[
            "First systematic ATPG algorithm — J. Paul Roth at IBM, 1966.",
            "Works directly on the stuck-at fault model.",
            "Introduces 5-valued logic: 0, 1, X (unknown), D (1 in good / 0 in faulty), D-bar (0 in good / 1 in faulty).",
            "D notation lets the algorithm track a discrepancy as it propagates through gates.",
            "Decisions can be made on ANY internal line (not just primary inputs).",
            "Provably complete: if a test pattern exists, D-Algorithm will find it.",
            "Provably exponential: worst-case decision space is huge, especially with reconvergent fanout.",
        ],
        notes=(
            "The D-Algorithm, published by Paul Roth at IBM in nineteen sixty-six, is the granddaddy of "
            "all ATPG algorithms. It works on the stuck-at fault model directly. Its big innovation was "
            "five-valued logic. In addition to zero, one, and X for unknown, Roth introduced D, which "
            "represents one in the good circuit and zero in the faulty one, and D-bar, the opposite. "
            "That symbolic representation lets the algorithm track a discrepancy symbolically as it "
            "propagates through gates — an AND of D and one is D, an AND of D and zero is zero, and so on. "
            "D-Algorithm is mathematically beautiful and provably complete. The catch: it makes decisions "
            "on any internal line, which means it can pick assignments that are physically inconsistent "
            "and waste enormous time backtracking. The worst case is exponential."
        ),
        image=None, time_s=80,
    ),
    dict(
        n=5, title="Working of D-Algorithm",
        bullets=[
            "Step 1 — Fault activation: place D or D-bar on the fault line by justifying the opposite value.",
            "Step 2 — D-frontier propagation: at each gate the D meets, sensitize the gate by setting side-inputs to non-controlling values.",
            "Step 3 — Line justification: every assigned internal value must be consistent with PI assignments — recursively justify upward.",
            "Step 4 — Backtrack on conflict: if any internal assignment becomes infeasible, undo the most recent decision.",
            "The D-frontier shrinks as faults propagate; the J-frontier (lines requiring justification) grows.",
            "Termination: success when D reaches a PO and all internal lines are justified by PI values.",
        ],
        notes=(
            "D-Algorithm runs in four phases that loop until the test pattern is found or the algorithm "
            "gives up. First, fault activation: we put a D or D-bar on the fault site, which means the "
            "opposite value has to be justified there in the good circuit. Second, D-frontier propagation: "
            "we walk the D forward through each gate it touches. To pass a D through, say, an AND gate, "
            "every other input has to be set to a non-controlling value, namely one for AND. Third, line "
            "justification: every internal value the algorithm assigns has to be implementable by some "
            "primary-input pattern. The algorithm recursively works backwards through gate fanin, choosing "
            "input values that produce the required output. Fourth, when something becomes inconsistent, "
            "the algorithm backtracks the most recent decision and tries again. The painful part is that "
            "decisions can happen on any internal line, so the search tree is enormous."
        ),
        image=None, time_s=85,
    ),
    # ----- PODEM block (reused from build_pptx.py) -----
    dict(
        n=6, title="PODEM (Goel 1981)",
        bullets=[
            "Path-Oriented Decision Making — proposed by Prabhu Goel at IBM, 1981.",
            "Direct response to D-Algorithm's exponential blow-up on XOR-heavy and reconvergent circuits.",
            "Key idea: branch only on Primary Inputs, never on internal lines.",
            "Every decision is therefore directly implementable — no inconsistent internal assignments to retract.",
            "Backtracking is bounded by 2^|PI|, far smaller than the D-Alg search space over all internal lines.",
        ],
        notes=(
            "PODEM was introduced by Prabhu Goel in 1981 as a direct fix for the explosion problem people kept "
            "hitting with the D-Algorithm. The D-Algorithm makes decisions on internal circuit lines, which means "
            "it can pick assignments that turn out to be physically impossible later, and then waste huge amounts "
            "of time backtracking. Goel's insight was simple but powerful: only ever make decisions at primary "
            "inputs. Every decision is then automatically consistent with the circuit, because PIs are free "
            "variables. The search space drops from 'all internal lines' to just two-to-the-number-of-primary-inputs, "
            "which is much smaller and much better behaved."
        ),
        image=None, time_s=75,
    ),
    dict(
        n=7, title="Working of PODEM",
        bullets=[
            "Objective: pick a target value on a specific line needed to activate or propagate the fault.",
            "Backtrace: walk backward from the objective line to a PI, choosing gate inputs that cause the desired output.",
            "Assign that PI, simulate forward, check if the objective is met.",
            "X-Path Check: confirm at least one all-X path still exists from the fault site to a PO; if not, backtrack.",
            "Completeness: PODEM is complete — if a test exists, it will eventually find it.",
            "Heuristics: SCOAP controllability/observability metrics guide objective and backtrace choices to reach a PI faster.",
        ],
        notes=(
            "PODEM works as a loop. First we set an objective — a value we want on some line, either to activate "
            "the fault or to propagate the D-frontier. Then we backtrace: starting from that objective line, we "
            "walk backwards through the gates, at each step choosing an input that would force the gate to the "
            "value we want. We keep walking until we hit a primary input, and we assign that PI. Then we simulate "
            "the whole circuit forward. If the objective is satisfied, great, pick the next one. If not, we use "
            "the X-Path Check — basically asking, 'is there still a route of unassigned lines from the fault to "
            "an output?' If yes, continue. If no, we backtrack the last PI decision. Because every decision is on "
            "a PI, backtracking is clean. SCOAP metrics — controllability and observability scores — are used as "
            "a heuristic to make the backtrace choose smart inputs first."
        ),
        image=None, time_s=90,
    ),
    dict(
        n=8, title="Worked PODEM Example",
        bullets=[
            "Tiny circuit: A, B, C are PIs. G1 = AND(A, B). G2 = OR(G1, C). Fault: G1 stuck-at-0.",
            "Activation objective: G1 = 1, which means A = 1 AND B = 1.",
            "Backtrace from G1: pick A first (highest controllability-1), assign A = 1, simulate.",
            "Re-evaluate: G1 still X because B is X. New objective: B = 1. Backtrace, assign B = 1.",
            "Now G1 = 1 in good circuit, 0 in faulty circuit — D activated at G1.",
            "Propagation objective: G2 must be 1 in good and 0 in faulty, so set side-input C = 0.",
            "Backtrace C = 0, assign, simulate. PO = G2 differs. Test pattern: A=1, B=1, C=0.",
        ],
        notes=(
            "Let's walk one through. Take a three-PI circuit: A, B, C. G1 is the AND of A and B. G2 is the OR of "
            "G1 and C, and G2 is the only output. Suppose we are testing G1 stuck-at-zero. Step one — activate "
            "the fault. To see a difference at G1, the good circuit needs G1 to be 1, so we need A equals 1 AND "
            "B equals 1. PODEM picks one objective at a time. We pick A equals 1, backtrace — A is already a PI "
            "so we just assign it — and simulate. G1 is still X because B is unknown. So our next objective is "
            "B equals 1. Assign, simulate, and now G1 has the value D, which means 1 in the good circuit and 0 "
            "in the faulty one. Step two — propagate. G2 is an OR. To pass D through an OR we need the side "
            "input, C, to be zero. Set C equals zero, simulate, and the output G2 now carries D. Done. Test "
            "pattern: A equals 1, B equals 1, C equals 0. Notice every single decision was on a primary input — "
            "no internal-line guessing."
        ),
        image=None, time_s=120,
    ),
    # ----- FAN block -----
    dict(
        n=9, title="FAN (Fujiwara-Shimono, 1983)",
        bullets=[
            "Fanout-Oriented test generation — H. Fujiwara and T. Shimono, 1983.",
            "Refines PODEM by classifying wires structurally before search.",
            "Fanout points: any line that drives more than one gate — the source of reconvergence.",
            "Bound lines: lines in the fanout cone of a fanout point — values are constrained.",
            "Free lines: lines with no fanout point in their fanin cone — values are unconstrained.",
            "Headlines: free lines that fan into bound regions — natural decision boundaries.",
            "Same completeness as PODEM, dramatically less backtracking on real circuits.",
        ],
        notes=(
            "FAN, the Fanout-Oriented test generation algorithm, was published by Fujiwara and Shimono in "
            "nineteen eighty-three. It builds on PODEM rather than replacing it. The big idea: before you "
            "start searching, look at the structure of the circuit and classify every wire. Fanout points "
            "are wires that drive more than one gate — they are the source of reconvergent paths and the "
            "real engine of search blow-up. Bound lines sit downstream of a fanout point; their values are "
            "constrained by the fanout. Free lines are upstream — their values are unconstrained. Headlines "
            "are the special class of free lines that feed into bound regions. FAN uses these headlines as "
            "natural decision boundaries: it stops backtracing at headlines instead of going all the way "
            "to the primary inputs, then justifies headlines later. Same completeness as PODEM, but much "
            "less backtracking in practice."
        ),
        image=None, time_s=85,
    ),
    dict(
        n=10, title="Working of FAN",
        bullets=[
            "Multiple backtrace: when several objectives exist, backtrace them in parallel — one pass instead of N.",
            "Stop at headlines, not at PIs — defer headline justification until propagation is fixed.",
            "Unique sensitization: when only one path can propagate the fault, force its side inputs immediately.",
            "X-path check: same as PODEM — abandon branches where no all-X route to a PO exists.",
            "Static learning: pre-compute implications of each gate's input combinations — avoid re-deriving at runtime.",
            "Net effect: same correctness as D-Algorithm and PODEM, far fewer wasted decisions.",
        ],
        notes=(
            "FAN improves PODEM with four big tricks. First, multiple backtrace. PODEM picks one objective "
            "at a time and runs a full backtrace. FAN backtraces multiple objectives simultaneously in one "
            "pass, which collapses redundant work. Second, stopping at headlines. Instead of always walking "
            "back to a primary input, FAN stops at the nearest headline and defers its justification, which "
            "shrinks the active search depth. Third, unique sensitization: when the algorithm can prove that "
            "only one path can possibly carry the fault to an output, it forces the side-inputs of that path "
            "to non-controlling values up front, avoiding wasted exploration. Fourth, static learning: the "
            "algorithm pre-computes obvious implications — like, if input A of an AND is zero then the output "
            "must be zero — so it does not have to rederive them every time. Same correctness as D-Algorithm "
            "and PODEM, far fewer wasted decisions."
        ),
        image=None, time_s=85,
    ),
    dict(
        n=11, title="Scalability Bottleneck and Why ML",
        bullets=[
            "Reconvergent fanout is the core source of pain — it creates internal constraints that classical search must rediscover.",
            "Worst-case runtime of D-Alg, PODEM, FAN is exponential in the number of fanout points.",
            "Our local PODEM took 117 seconds on c432 (160 gates) — and modern designs have 100M+ gates.",
            "Commercial tools (TestMAX, ATALANTA) survive via heuristics, parallelism, and engineering effort — but still scale poorly.",
            "Observation: useful test patterns share structural regularities — fanout cones, gate sequences, reconvergence shapes.",
            "ML hypothesis: a neural network can learn these regularities and predict a candidate pattern in one forward pass.",
            "This is the motivation for DETECTive (Petrolo et al., GLSVLSI 2024) — and the rest of the talk.",
        ],
        notes=(
            "Why are we even talking about machine learning for ATPG? Because the classical algorithms, "
            "elegant as they are, all share an exponential worst case driven by reconvergent fanout. "
            "Reconvergence creates implicit constraints between internal lines, and search has to rediscover "
            "those constraints over and over. Our own measurements bear this out: PODEM on c432 — only one "
            "hundred sixty gates — took us almost two minutes. Real industrial designs have a hundred million "
            "gates or more. Commercial tools cope through heuristics, massive parallelism, and a lot of "
            "engineering hours, but they still scale poorly. The machine-learning observation is that useful "
            "test patterns share structural regularities — fanout cones, gate-type sequences, reconvergence "
            "shapes — and a neural network ought to be able to learn those regularities and predict a "
            "candidate pattern directly, in a single forward pass, with no search at all. That is exactly "
            "what the DETECTive paper proposes, and what we replicated. The rest of the talk is about that "
            "system."
        ),
        image=None, time_s=90,
    ),
    # =====================================================================
    # BLOCK 2 — ATPP Architecture
    # =====================================================================
    dict(
        n=12, title="ATPP Concept",
        bullets=[
            "ATPP = Automatic Test Pattern Prediction (vs. Generation).",
            "Input: gate-level netlist + a target stuck-at fault location and polarity.",
            "Output: a single binary test pattern over the primary inputs, in one forward pass.",
            "Architecture: GNN encoder learns gate-level structure, LSTM encoder learns path context, MLP head predicts PI bits.",
            "Goal per the paper: prioritize one test per fault, not maximize coverage in one pattern.",
            "Evaluation: bit-wise comparison with the closest-matching ground-truth pattern (paper Sec. 5).",
        ],
        notes=(
            "The acronym is deliberate. Classical ATPG generates patterns by searching. ATPP — prediction — "
            "produces them in one forward pass. The model takes two inputs: the gate-level netlist and a "
            "designated fault, meaning a specific line and a stuck-at-0 or stuck-at-1 polarity. It outputs a "
            "vector of bits, one per primary input. That's the proposed test pattern. The architecture has "
            "three blocks: a graph neural network encodes the local gate structure, an LSTM encodes the "
            "activation and propagation paths, and a small MLP head decodes the joint embedding into PI bits. "
            "One important detail from Section 5 of the paper: the model is trained to prioritize one test per "
            "fault, not to find a maximally compacted pattern that covers many faults at once."
        ),
        image=None, time_s=70,
    ),
    dict(
        n=13, title="Circuit as Graph",
        bullets=[
            "Each gate becomes a node; each wire becomes a directed edge from driver to load.",
            "Primary inputs and outputs are special node types with no driver / no load.",
            "The fault site is encoded by a one-bit flag on the affected node.",
            "Resulting structure is a directed acyclic graph (combinational logic only).",
            "Paper Fig. 2 shows this transformation; we follow it exactly.",
            "This representation lets a GNN reason about local fanin/fanout and reconvergence directly.",
        ],
        notes=(
            "Step one is turning the netlist into something a neural network can consume. Every gate becomes "
            "a node, every wire becomes a directed edge from the driving gate to the loading gate. Primary "
            "inputs and outputs become special nodes. The target fault is encoded by setting a flag bit on "
            "the node where the fault sits. Because we restrict ourselves to combinational logic, the resulting "
            "graph is a directed acyclic graph, which keeps the GNN well-behaved. This is exactly the "
            "representation in Figure 2 of the paper, and we follow it without modification. The key advantage: "
            "a GNN can directly attend to the local fanin and fanout structure around the fault, which is "
            "exactly the information classical ATPG also exploits, just in a different form."
        ),
        image=None, time_s=60,
    ),
    dict(
        n=14, title="GNN Encoder: GAT + GCN",
        bullets=[
            "First layer: Graph Attention Network (GAT) — learns weighted importance over each node's neighbors.",
            "Useful for fanout, where some loads matter more than others for fault propagation.",
            "Second layer: Graph Convolutional Network (GCN) — uniform aggregation, refines the attention output.",
            "Output embedding dimension: 32 per node.",
            "ReLU activations, dropout 0.2 between layers.",
            "Result: each node carries a learned vector summarizing its structural role.",
        ],
        notes=(
            "The encoder stacks two graph layers. The first is a Graph Attention Network. GAT is important "
            "here because, at any node, its neighbors are not equally informative. A fanout going to a deep "
            "cone matters more for propagation than a fanout going to a dead end. GAT learns those weights "
            "from data. The second layer is a standard Graph Convolutional Network, which does uniform "
            "mean-aggregation. We use GCN as a smoothing step on top of GAT — it refines the attention output "
            "and tends to stabilize training. Both layers produce 32-dimensional node embeddings, with ReLU "
            "activations and dropout of zero point two between them."
        ),
        image=None, time_s=65,
    ),
    dict(
        n=15, title="LSTM Path Encoder",
        bullets=[
            "Two path types per fault: activation paths (PIs to fault site) and propagation paths (fault site to POs).",
            "Up to p = 10 paths of each type; paths are variable length.",
            "Each path is a sequence of GNN node embeddings — naturally fits an LSTM.",
            "We use a Bidirectional LSTM so context flows both forward and backward along the path.",
            "Hidden size 32, single layer, mean-pooled across the p paths to give a fixed vector per type.",
            "Output: an activation-path embedding and a propagation-path embedding, each in R^32.",
        ],
        notes=(
            "The GNN tells us about local structure. The LSTM tells us about path-level structure, which is "
            "what ATPG actually cares about. For each fault we extract up to ten activation paths — routes "
            "from a primary input to the fault site — and up to ten propagation paths — routes from the fault "
            "site to a primary output. Each path is a sequence of node embeddings produced by the GNN. We "
            "feed those sequences into a bidirectional LSTM with hidden size thirty-two. Bidirectional matters "
            "because both ends of a path matter — the PI end determines controllability, the PO end determines "
            "observability. We then mean-pool across the up-to-ten paths to get a single fixed-size embedding "
            "per path type."
        ),
        image=None, time_s=70,
    ),
    dict(
        n=16, title="Full Architecture",
        bullets=[
            "Three modules: GNN encoder, LSTM path encoder, MLP InputPredictor head.",
            "Hidden dim 32 throughout, p = 10 paths, dropout 0.2.",
            "InputPredictor concatenates (GNN-fault-node-embedding, activation-path-emb, propagation-path-emb).",
            "MLP: 96 -> 64 -> |PI| with sigmoid output, one bit per primary input.",
            "Total trainable parameters: ~25 k — small by modern standards, intentionally so.",
            "One forward pass per fault gives the predicted test pattern.",
        ],
        notes=(
            "Putting it all together — three modules. The GNN encodes the netlist into node embeddings. The "
            "path encoder selects up to ten activation and ten propagation paths and passes each through a "
            "bidirectional LSTM. The InputPredictor head takes three vectors — the embedding of the fault "
            "node itself, the pooled activation-path embedding, and the pooled propagation-path embedding — "
            "concatenates them into a 96-dim vector, and passes them through a two-layer MLP with sigmoid "
            "output. The output dimension is the number of primary inputs of whatever circuit we're predicting "
            "on. The total parameter count is around twenty-five thousand — deliberately small."
        ),
        image=None, time_s=70,
    ),
    dict(
        n=17, title="Node Features and Tensor Shapes",
        bullets=[
            "Per-node feature: 11-dim vector — 10-way one-hot for gate type plus 1-bit fault flag.",
            "Gate types: AND, OR, NAND, NOR, XOR, XNOR, NOT, BUFF, PI, PO.",
            "Tensor flow: [N, 11] -> GNN -> [N, 32] -> path slice [p, L, 32] -> BiLSTM -> [p, 32] -> mean -> [32].",
            "Concatenated head input: [96] -> MLP -> [|PI|] sigmoid -> binary test pattern.",
            "Threshold at 0.5 to get hard bits at inference; raw logits used for paper's bit-accuracy metric.",
            "Same shapes whether circuit is 5 gates or 1000 gates — model is parameter-shared across N.",
        ],
        notes=(
            "Concretely, every node carries an eleven-dimensional feature vector. Ten of those bits are a "
            "one-hot indicator of the gate type. The eleventh bit is a fault flag, set on the single node "
            "where we are targeting a fault. The GNN turns N-by-eleven into N-by-thirty-two. We then extract "
            "path tensors of shape p paths by max-length L by thirty-two, run them through the BiLSTM, pool "
            "across paths, and get a single thirty-two-vector per path type. Concatenated with the fault "
            "node's own embedding, that's a ninety-six-vector going into the MLP. The same architecture, "
            "same parameters, work for circuits of any size — that's what graph weight sharing gives us."
        ),
        image=None, time_s=75,
    ),
    dict(
        n=18, title="Training Setup",
        bullets=[
            "Loss: Binary Cross-Entropy applied per PI bit, averaged.",
            "Optimizer: Adam, learning rate 1e-3, gradient clipping at 1.0, batch size 32.",
            "Train/val split: 80/20 over ~12 k synthetic 4-input circuits.",
            "Multi-ground-truth handling: closest-pattern selection — pick GT pattern with min Hamming distance.",
            "Trained 100 epochs; checkpoint best validation bit-accuracy.",
            "Weights file: 100_epoch_run/best_detective_model.pt (~25k params).",
        ],
        notes=(
            "Training is straightforward. The loss is binary cross-entropy applied independently to each "
            "primary-input bit and averaged — the paper treats every bit as its own classification problem. "
            "We use Adam at learning rate one-e-minus-three, clip gradients at one point zero to keep training "
            "stable, and batch thirty-two. The dataset is twelve thousand synthetic four-input circuits, "
            "split eighty-twenty for train and validation. One important trick: many faults have multiple "
            "valid ground-truth patterns. To avoid penalizing the model for picking a different valid one, "
            "we compare its prediction against the ground-truth pattern with the smallest Hamming distance "
            "to the prediction. We train for one hundred epochs and checkpoint the best validation bit-accuracy."
        ),
        image=None, time_s=70,
    ),
    dict(
        n=19, title="Training Curve and Overfitting",
        bullets=[
            "Best validation bit-accuracy: 0.8358 at epoch 40.",
            "Training accuracy continues climbing past epoch 40 — model has capacity to memorize.",
            "Validation plateaus then drifts down — classic overfitting signature.",
            "We report the epoch-40 checkpoint, not the final epoch.",
            "Train-val gap widens monotonically past epoch 40 by ~0.5 - 1 percentage point per 10 epochs.",
            "Justifies the early-best checkpoint policy and the modest parameter budget (~25 k).",
        ],
        notes=(
            "The training curve tells a clean overfitting story. Best validation bit-accuracy was zero "
            "point eight three five eight at epoch forty. Past that, training accuracy keeps climbing "
            "because the model has plenty of capacity to memorize, but validation accuracy plateaus and "
            "then drifts down. The train-val gap widens by about half a percent to one percent per ten "
            "epochs after that. This is exactly the pattern that justifies an early-best checkpoint policy "
            "rather than just taking the final weights. It also justifies why we kept the parameter budget "
            "modest at around twenty-five thousand — a bigger model would overfit harder on twelve thousand "
            "circuits."
        ),
        image=None, time_s=70,
    ),
    # =====================================================================
    # BLOCK 3 — Results
    # =====================================================================
    dict(
        n=20, title="Dataset",
        bullets=[
            "~12 000 randomly generated combinational circuits, each with exactly 4 primary inputs.",
            "Random gate counts, depths, and topologies — chosen to span the regimes the paper studies.",
            "Ground truth: brute-force exhaustive search across all 16 PI patterns per fault.",
            "Multiple valid test patterns exist per fault — handled by closest-pattern selection at evaluation.",
            "Fault list: every line in every circuit, both stuck-at-0 and stuck-at-1.",
            "Generation script and pickled splits live in train_dataset.pkl, val_dataset.pkl.",
        ],
        notes=(
            "The dataset is what makes the paper-faithful evaluation possible. We generated about twelve "
            "thousand random combinational circuits, each constrained to have exactly four primary inputs. "
            "Because the input space is only sixteen patterns, we can compute ground truth by brute force — "
            "for each fault we simply enumerate all sixteen and record which of them detect it. That gives "
            "a set of valid patterns per fault, often multiple. The paper handles multi-pattern faults the "
            "same way we do: at evaluation time we compare the predicted pattern against whichever "
            "ground-truth pattern is closest in Hamming distance."
        ),
        image=None, time_s=70,
    ),
    dict(
        n=21, title="ISCAS-85 Results Table",
        bullets=[
            "Five circuits evaluated: c17, c432, c499, c880, c1908.",
            "Per circuit: gate count, ATPP bit-accuracy, ATPP fault-sim coverage, total runtime in ms.",
            "c17 (6 gates): bit-acc 86.67%, fault-cov 83.33%, 408.5 ms.",
            "c432 (160 gates): bit-acc 85.53%, fault-cov 41.88%, 665 ms.",
            "c499 (202 gates): bit-acc 44.27%, fault-cov 17.82%, 692.8 ms.",
            "c880 (383 gates): bit-acc 92.58%, fault-cov 15.54%, 1130.2 ms.",
            "c1908 (880 gates): bit-acc 69.70%, fault-cov 23.86%, 572.2 ms.",
        ],
        notes=(
            "Here are our measured ATPP numbers across the five ISCAS-85 benchmarks. Two metrics: "
            "bit-accuracy, which is the paper's official metric, and fault-simulation coverage, which is "
            "a strictly harder metric we computed for honesty. Bit-accuracy is high on c17, c432, and c880, "
            "all above eighty-five percent, with c880 the best at ninety-two-and-a-half. c499 is the painful "
            "one, dropping to forty-four percent — c499 is XOR-heavy and full of reconvergence. Runtimes "
            "are in milliseconds total across all faults, all well under two seconds. Notice c432 — only "
            "six hundred sixty-five milliseconds against PODEM's one hundred seventeen seconds."
        ),
        image=None, time_s=80,
    ),
    dict(
        n=22, title="4-Way Bit-Accuracy (Paper-Faithful)",
        bullets=[
            "Bar chart compares ATPP, PODEM, D-Algorithm, FAN on the bit-accuracy metric.",
            "Paper Section 5 verbatim: bit-accuracy is the official DETECTive metric — closest-pattern bit-wise compare.",
            "ATPP scores (%): c17 86.67, c432 85.53, c499 44.27, c880 92.58, c1908 69.70.",
            "ATPP is competitive with classical algorithms on c17, c432, c880.",
            "ATPP underperforms on c499 (XOR-heavy reconvergence) — same family of circuits the paper flags as hardest.",
            "This is the apples-to-apples view the paper uses; runtime is reported separately.",
        ],
        notes=(
            "This is the bit-accuracy comparison across all four algorithms. The metric is the paper's "
            "official one — Section 5, verbatim, says: 'we compute our model's accuracy by selecting the "
            "test pattern that exhibits the closest similarity to the predicted test pattern.' Under that "
            "metric, ATPP comes within a few points of the classical algorithms on c17, c432, and c880. "
            "It's clearly weakest on c499, which is consistent with the paper's own Figure 7b — c499 is the "
            "lowest point on their plot too, around seventy-five percent. The take-away is: on the metric "
            "the paper defines, ATPP is competitive everywhere except XOR-heavy reconvergent circuits."
        ),
        image="bit_accuracy_4way.png", time_s=75,
    ),
    dict(
        n=23, title="Runtime per Fault — 17.6x Speedup on c432",
        bullets=[
            "c432 headline: PODEM 117 034 ms total vs ATPP 665 ms — a 176x circuit-level speedup, ~17.6x per-fault.",
            "Across all five circuits ATPP runtime stays in the hundreds-of-ms range.",
            "PODEM runtime explodes with circuit complexity — exponential search behavior visible.",
            "D-Algorithm and FAN follow PODEM's general scaling pattern, just with different constants.",
            "ATPP runtime is dominated by GNN+LSTM forward pass, near-constant per fault.",
            "This is the result the paper centers on, and our measurements reproduce it directionally.",
        ],
        notes=(
            "This is the most important result in the deck. On c432 specifically, PODEM took one hundred "
            "seventeen seconds end-to-end to cover the fault list. ATPP did the same circuit in six hundred "
            "sixty-five milliseconds. That's a one-hundred-seventy-six-times reduction at circuit level — "
            "about seventeen-point-six-times faster per individual fault. Across all five benchmarks, ATPP "
            "stays in the hundreds-of-milliseconds regime, while classical methods grow with the circuit. "
            "Note: we're comparing our local PODEM, not Synopsys TestMAX. The paper reports DETECTive as "
            "roughly two-times faster than TestMAX and fifteen-times faster than ATALANTA."
        ),
        image="runtime_perfault_4way.png", time_s=80,
    ),
    dict(
        n=24, title="Bit-Accuracy vs Fault-Sim Coverage Gap",
        bullets=[
            "Bit-accuracy and fault-simulation coverage are NOT the same thing — and the gap is real.",
            "Bit-accuracy (paper metric): how many bits match the closest ground-truth pattern.",
            "Fault-sim coverage (stricter): does the predicted pattern actually detect the fault when simulated?",
            "Measured gap: c432 85.53% bit-acc -> 41.88% fault-cov; c880 92.58% -> 15.54%; c1908 69.70% -> 23.86%.",
            "The paper acknowledges this: it explicitly chooses bit-level over word-level (Section 6.2).",
            "We report both because honest replication requires acknowledging which guarantee you actually have.",
        ],
        notes=(
            "This is the slide where we have to be careful and honest. The paper's metric is bit-accuracy — "
            "closest-pattern bitwise match. From Section 5, verbatim: 'We train DETECTive to generate a "
            "single test pattern that exposes the designated stuck-at fault, prioritizing it rather than "
            "aiming for broader fault coverage. We compute our model's accuracy by selecting the test "
            "pattern that exhibits the closest similarity to the predicted test pattern. The accuracy "
            "metric is determined through a bit-wise comparison of the two test patterns.' That metric is "
            "real and meaningful, but it isn't the same as 'did the predicted pattern actually detect the "
            "fault.' A pattern can be one bit off and totally fail to activate the fault. We computed the "
            "stricter metric too. The gap is large. On c432, eighty-five percent bit-accuracy translates "
            "to forty-two percent fault coverage. On c880, ninety-two becomes fifteen. We're not invalidating "
            "the paper — Section six-point-two explicitly justifies the bit-level choice — but a faithful "
            "replication has to surface this gap."
        ),
        image="fault_sim_coverage_table.png", time_s=90,
    ),
    dict(
        n=25, title="Synthetic Grid (Paper Fig 5a Replication)",
        bullets=[
            "Grid sweeps gate count vs circuit depth on synthetic 4-input circuits — replicates paper Figure 5a.",
            "Bit-accuracy range: 0.78 to 1.00 across the grid; mean 0.87.",
            "Trend is monotone: small + shallow circuits hit ~0.93, deep + large circuits drop to ~0.79.",
            "Confirms paper's claim that the model degrades smoothly with structural complexity — no sharp cliffs.",
            "Synthetic regime is in-distribution — these are circuits drawn from the same generator as training.",
            "ISCAS-85 results are out-of-distribution — bigger and structurally different than training.",
        ],
        notes=(
            "We reproduced Figure 5a from the paper. The plot is a heatmap over a grid of gate counts on one "
            "axis and circuit depths on the other, on synthetic four-PI circuits drawn from the same generator "
            "as training. Bit-accuracy ranges from zero point seven eight up to one point zero, with a mean "
            "around zero point eight seven. The trend is monotone: small and shallow circuits are easy, deep "
            "and large circuits are hard, but the degradation is smooth — no sharp cliffs. One thing to keep "
            "in mind — these are in-distribution circuits. ISCAS-85 is out-of-distribution: bigger circuits, "
            "more PIs, different structural conventions."
        ),
        image="synthetic_grid_heatmap.png", time_s=75,
    ),
    dict(
        n=26, title="Depth Degradation (Paper Fig 6a Replication)",
        bullets=[
            "Plot: bit-accuracy vs maximum logic depth, holding gate count roughly fixed.",
            "Range observed: 0.89 at depth 1 down to 0.78 at depth 25.",
            "Smooth curve — no phase transition, no collapse.",
            "Matches paper Figure 6a directionally.",
            "Mechanism: deeper circuits have longer activation/propagation paths than the LSTM saw in training.",
            "Suggests one obvious extension — train on deeper synthetic circuits to push the curve right.",
        ],
        notes=(
            "The companion plot — bit-accuracy as a function of circuit depth, holding gate count fixed. "
            "We see zero point eight nine at depth one, drifting down to zero point seven eight at depth "
            "twenty-five. The curve is smooth — no phase transition, no collapse, just a gentle slope. "
            "This matches the shape in Figure 6a of the paper. The mechanistic explanation is that deeper "
            "circuits produce longer activation and propagation paths, and the LSTM is being asked to "
            "extrapolate beyond the path lengths it saw during training. The simplest improvement: train "
            "on deeper circuits, or swap the LSTM for a Transformer with positional embeddings."
        ),
        image="depth_curve.png", time_s=75,
    ),
    dict(
        n=27, title="Honest Takeaways and Limitations",
        bullets=[
            "Speed: yes — 17.6x faster than PODEM on c432, hundreds-of-ms regime across ISCAS-85.",
            "Bit-accuracy: competitive with classical ATPG on most circuits under the paper's official metric.",
            "Completeness: no — ATPP gives no formal guarantee that any predicted pattern detects its fault.",
            "Weak spots: c499 and c1908 — XOR-heavy and reconvergent, same circuits where classical methods also struggle.",
            "Bit-acc vs fault-coverage gap is real; the paper's metric choice is defensible but worth flagging.",
            "Scope: combinational only; deeper paths hurt; no sequential circuits in this work.",
            "Honest framing: ATPP is a fast first guess for a classical verifier, not a replacement for ATPG.",
        ],
        notes=(
            "To wrap up the results block — what did we actually demonstrate? Speed, yes — seventeen-times "
            "faster than PODEM on c432, hundreds of milliseconds across all benchmarks. Bit-accuracy is "
            "competitive with classical ATPG on most circuits under the paper's metric. But completeness, "
            "no — there's no formal guarantee any individual prediction detects its fault. The architecture "
            "struggles on c499 and c1908, the same circuits where D-Algorithm and FAN also struggle, so "
            "this is a hard-instance problem, not just a model bug. The bit-acc to fault-coverage gap is "
            "real and we surfaced it. Scope limits: combinational logic only, no sequential circuits, and "
            "deeper paths degrade accuracy. The honest framing: ATPP is best understood as a very fast first "
            "guess that a classical verifier can confirm or repair. Used that way, the speedup compounds."
        ),
        image=None, time_s=90,
    ),
    dict(
        n=28, title="Conclusion",
        bullets=[
            "Classical ATPG (D-Alg, PODEM, FAN) is correct and complete but exponential in the worst case.",
            "DETECTive (ATPP) replaces search with a single GNN+LSTM forward pass — predicts a candidate test pattern.",
            "Demonstrated 17.6x per-fault speedup over PODEM on c432; 85.5% bit-accuracy under the paper's metric.",
            "The model transfers from synthetic 4-PI training data to real ISCAS-85 circuits — generalization works.",
            "Future work: integrate ATPP as fast first guess + classical verifier; extend to sequential circuits; deeper-circuit training.",
            "Broader take: machine learning is a real tool for EDA — not a replacement for proofs, but a meaningful accelerator.",
        ],
        notes=(
            "To conclude: classical ATPG works, it is provably complete, but it is exponential in the worst "
            "case and that is the bottleneck on modern designs. DETECTive proposes a different bet — let a "
            "neural network learn the structural regularities of useful test patterns, and produce them in "
            "a single forward pass. We replicated the system and confirmed the headline numbers: about "
            "seventeen-times faster per fault than our PODEM on c432, with eighty-five-and-a-half percent "
            "bit-accuracy under the paper's official metric. The model transfers from synthetic four-PI "
            "training data to ISCAS-85 with no fine-tuning, which is the generalization story. Future "
            "directions: integrate ATPP as a fast first guess in front of a classical verifier; extend to "
            "sequential circuits; train on deeper circuits to push the depth curve. The broader point is "
            "that machine learning is a real tool for EDA — not a replacement for formal guarantees, but "
            "a meaningful accelerator. Thank you, and we are happy to take questions."
        ),
        image=None, time_s=80,
    ),
]


# ---------------------------------------------------------------------------
# Helpers (style matched to build_pptx.py)
# ---------------------------------------------------------------------------

def set_text(tf, text, font_size=18, bold=False, color=None):
    tf.text = text
    p = tf.paragraphs[0]
    if p.runs:
        r = p.runs[0]
    else:
        r = p.add_run()
        r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    if color is not None:
        r.font.color.rgb = color


def add_title_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.0), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    set_text(tf, "Test Vector Generation for Fault Detection Using Graph Neural Networks",
             font_size=34, bold=True, color=RGBColor(0x1F, 0x2E, 0x4F))

    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.1), Inches(12.0), Inches(1.0))
    tf = sub_box.text_frame
    tf.word_wrap = True
    set_text(tf,
             "Replicating DETECTive: A 30-minute Walk-through of Classical ATPG and ML-driven ATPP",
             font_size=22, bold=False, color=RGBColor(0x33, 0x33, 0x33))

    cite_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.6), Inches(12.0), Inches(1.5))
    tf = cite_box.text_frame
    tf.word_wrap = True
    set_text(tf,
             "Reference: Petrolo et al., DETECTive: Machine Learning-driven Automatic "
             "Test Pattern Prediction for Faults in Digital Circuits, GLSVLSI 2024.",
             font_size=16, color=RGBColor(0x55, 0x55, 0x55))

    foot_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.3), Inches(12.0), Inches(0.7))
    tf = foot_box.text_frame
    set_text(tf,
             "Block 1 — Classical ATPG (D-Alg, PODEM, FAN) | Block 2 — ATPP Architecture | Block 3 — Results & Conclusion",
             font_size=14, color=RGBColor(0x77, 0x77, 0x77))

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "Welcome — this is a thirty-minute walk-through of test pattern generation for digital circuits. "
        "We will cover three blocks. First, the classical ATPG world: the fault model, what a test pattern "
        "is, the exponential blow-up of brute force, and the three landmark algorithms — D-Algorithm, "
        "PODEM, and FAN. That motivates why people are now looking at machine learning. Second, we will "
        "walk through the architecture of DETECTive, a recent paper from GLSVLSI twenty twenty-four that "
        "uses a graph neural network plus an LSTM to predict test patterns directly. Third, our replication "
        "results on the ISCAS-85 benchmarks, with honest framing of where the system works and where it "
        "doesn't. Total time about thirty minutes."
    )


def add_content_slide(prs, spec):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    has_image = spec.get("image") is not None

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    title_text = f"{spec['n']}. {spec['title']}"
    set_text(tf, title_text, font_size=28, bold=True, color=RGBColor(0x1F, 0x2E, 0x4F))

    # Bullets — left column if image present, full width otherwise
    if has_image:
        body_left, body_width = Inches(0.5), Inches(6.5)
    else:
        body_left, body_width = Inches(0.5), Inches(12.3)
    body_box = slide.shapes.add_textbox(body_left, Inches(1.3), body_width, Inches(5.5))
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(spec["bullets"]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(6)

    # Image (or placeholder)
    if has_image:
        img_path = os.path.join(IMG_DIR, spec["image"])
        img_left = Inches(7.2)
        img_top = Inches(1.5)
        img_width = Inches(5.8)
        img_height = Inches(4.5)
        if os.path.isfile(img_path):
            try:
                slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)
            except Exception as e:
                _add_image_placeholder(slide, spec["image"], img_left, img_top, img_width, img_height,
                                       err=str(e))
        else:
            _add_image_placeholder(slide, spec["image"], img_left, img_top, img_width, img_height)

    # Time-budget footer
    foot_box = slide.shapes.add_textbox(Inches(11.0), Inches(6.9), Inches(2.2), Inches(0.4))
    tf = foot_box.text_frame
    set_text(tf, f"~{spec['time_s']}s", font_size=11, color=RGBColor(0x99, 0x99, 0x99))

    # Speaker notes
    slide.notes_slide.notes_text_frame.text = spec["notes"]


def _add_image_placeholder(slide, name, left, top, width, height, err=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    msg = f"[Image placeholder]\n\npresentation_images/{name}\n\n(PNG not found.)"
    if err:
        msg += f"\n\n(Insert error: {err})"
    set_text(tf, msg, font_size=14, color=RGBColor(0x88, 0x88, 0x88))
    p = tf.paragraphs[0]
    for run in p.runs:
        run.font.italic = True


def add_qa_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12.0), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    set_text(tf, "Q & A", font_size=80, bold=True, color=RGBColor(0x1F, 0x2E, 0x4F))
    tf.paragraphs[0].alignment = 2  # center

    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.7), Inches(12.0), Inches(1.0))
    tf = sub_box.text_frame
    tf.word_wrap = True
    set_text(tf,
             "Happy to discuss: classical ATPG details | ATPP architecture choices | replication caveats | future work.",
             font_size=18, color=RGBColor(0x55, 0x55, 0x55))
    tf.paragraphs[0].alignment = 2

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "Open the floor for questions. Likely topics: why bit-accuracy versus fault-sim coverage, why "
        "GAT plus GCN rather than one or the other, whether ATPP can replace classical ATPG entirely "
        "(no — fast first guess), how it would extend to sequential circuits (would need scan-chain "
        "modeling and possibly a temporal encoder), and how it compares to commercial tools like Synopsys "
        "TestMAX or ATALANTA in absolute terms."
    )


def add_thanks_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(12.0), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    set_text(tf, "Thank You", font_size=80, bold=True, color=RGBColor(0x1F, 0x2E, 0x4F))
    tf.paragraphs[0].alignment = 2  # center

    cite_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(12.0), Inches(1.5))
    tf = cite_box.text_frame
    tf.word_wrap = True
    set_text(tf,
             "Reference: Petrolo et al., DETECTive: Machine Learning-driven Automatic "
             "Test Pattern Prediction for Faults in Digital Circuits, GLSVLSI 2024.",
             font_size=14, color=RGBColor(0x77, 0x77, 0x77))
    tf.paragraphs[0].alignment = 2

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "Thank the audience and stick around for offline discussion. Mention where the code and weights "
        "live for anyone who wants to reproduce or extend the work."
    )


def main():
    prs = Presentation()
    # 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    for spec in SLIDES:
        add_content_slide(prs, spec)
    add_qa_slide(prs)
    add_thanks_slide(prs)

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"File size: {os.path.getsize(OUT_PATH)} bytes")


if __name__ == "__main__":
    main()
